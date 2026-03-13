import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def create_presentation():
    prs = Presentation()

    # Logo path
    logo_path = r"C:\Users\katte\.gemini\antigravity\brain\e0228e8e-4f7d-4f20-a737-b36b4f4826c1\sortify_logo_mockup_1773286949141.png"

    # --- Slide 1: Title ---
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "Sortify"
    subtitle.text = "Reclaim Your Inbox: Focus on What Matters"

    if os.path.exists(logo_path):
        slide.shapes.add_picture(logo_path, Inches(3.5), Inches(4.5), width=Inches(3))

    # --- Slide 2: The Problem ---
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "The Problem: Email Overload"
    
    tf = slide.placeholders[1].text_frame
    tf.text = "Inbox noise is a productivity killer."
    p = tf.add_paragraph()
    p.text = "• Mundane 'Important' emails hide truly urgent ones."
    p = tf.add_paragraph()
    p.text = "• Constant notifications lead to focus fragmentation."
    p = tf.add_paragraph()
    p.text = "• Existing filters are often too broad or too complex."

    # --- Slide 3: The Solution ---
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "The Solution: Sortify"
    
    tf = slide.placeholders[1].text_frame
    tf.text = "A sleek, zero-cost urgent email highlighting tool."
    p = tf.add_paragraph()
    p.text = "• Hybrid rules-based engine for high precision."
    p = tf.add_paragraph()
    p.text = "• Local execution: High privacy, zero API costs."
    p = tf.add_paragraph()
    p.text = "• Instant urgency evaluation without heavy ML."

    # --- Slide 4: Key Features ---
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Key Features"
    
    tf = slide.placeholders[1].text_frame
    p = tf.add_paragraph()
    p.text = "• Android Home Screen Widget for at-a-glance urgency."
    p = tf.add_paragraph()
    p.text = "• Gmail & Outlook Integration via secure OAuth."
    p = tf.add_paragraph()
    p.text = "• Privacy-First: No full email bodies stored remotely."
    p = tf.add_paragraph()
    p.text = "• Custom Glassmorphism UI for a premium experience."

    # --- Slide 5: Tech Stack & Privacy ---
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Built for Stability & Privacy"
    
    tf = slide.placeholders[1].text_frame
    p = tf.add_paragraph()
    p.text = "• Backend: Spring Boot (Java 17) with H2 Database."
    p = tf.add_paragraph()
    p.text = "• Frontend: React Native for Android."
    p = tf.add_paragraph()
    p.text = "• Data Security: Local persistence for email snippets only."
    p = tf.add_paragraph()
    p.text = "• Modern Architecture: Native AppWidgetProvider module."

    # Save presentation
    output_path = r"C:\Users\katte\sortify\Sortify_Presentation.pptx"
    prs.save(output_path)
    print(f"Presentation saved to {output_path}")

if __name__ == "__main__":
    create_presentation()
